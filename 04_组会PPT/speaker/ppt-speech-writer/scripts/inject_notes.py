#!/usr/bin/env python3
"""Inject speaker notes into a .pptx file's notes pane."""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def ensure_pptx() -> None:
    try:
        import pptx  # noqa: F401
    except ImportError:
        sys.stderr.write("python-pptx is required. Install it in the active environment.\n")
        raise SystemExit(2)


def set_notes_text(slide, text: str) -> None:
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    lines = text.split("\n")
    tf.text = lines[0] if lines else ""
    for extra in lines[1:]:
        p = tf.add_paragraph()
        p.text = extra


def existing_notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return (slide.notes_slide.notes_text_frame.text or "").strip()


def parse_notes_file(path: Path) -> dict[int, str]:
    raw = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(raw, list):
        raise ValueError("notes.json must be a JSON array")
    notes: dict[int, str] = {}
    for entry in raw:
        if not isinstance(entry, dict) or "slide" not in entry or "notes" not in entry:
            raise ValueError(f"bad entry: {entry!r}")
        slide = int(entry["slide"])
        if slide in notes:
            raise ValueError(f"duplicate slide entry: {slide}")
        notes[slide] = str(entry["notes"])
    return notes


def inject(input_path: Path, output_path: Path, notes_map: dict[int, str], mode: str) -> list[str]:
    from pptx import Presentation

    prs = Presentation(str(input_path))
    total = len(prs.slides)
    log: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        if idx not in notes_map:
            log.append(f"slide {idx}: skipped (no entry in notes.json)")
            continue
        new_text = notes_map[idx].strip()
        existing = existing_notes_text(slide)
        if mode == "skip-if-present" and existing:
            log.append(f"slide {idx}: skipped (existing notes present)")
            continue
        if mode == "append" and existing:
            set_notes_text(slide, f"{existing}\n\n{new_text}")
            log.append(f"slide {idx}: appended ({len(new_text)} chars)")
        else:
            set_notes_text(slide, new_text)
            log.append(f"slide {idx}: replaced ({len(new_text)} chars)")

    for extra in sorted(k for k in notes_map if k < 1 or k > total):
        log.append(f"WARNING: notes.json entry for slide {extra} ignored (deck has {total} slides)")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    log.append(f"saved: {output_path}")
    return log


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description="Inject speaker notes into a .pptx file.")
    parser.add_argument("--input", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--notes", required=True, type=Path)
    parser.add_argument("--mode", choices=["replace", "append", "skip-if-present"], default="replace")
    args = parser.parse_args(argv)

    if not args.input.exists():
        sys.stderr.write(f"Input not found: {args.input}\n")
        return 1
    if not args.notes.exists():
        sys.stderr.write(f"Notes file not found: {args.notes}\n")
        return 1

    ensure_pptx()
    try:
        notes_map = parse_notes_file(args.notes)
    except (ValueError, json.JSONDecodeError) as exc:
        sys.stderr.write(f"Failed to parse notes file: {exc}\n")
        return 1
    for line in inject(args.input, args.output, notes_map, args.mode):
        print(line)
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
