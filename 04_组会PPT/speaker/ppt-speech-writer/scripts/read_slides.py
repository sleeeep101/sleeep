#!/usr/bin/env python3
"""Extract structured evidence from a .pptx deck.

The output is a JSON object with one entry per slide. It captures text frames,
tables, chart XML, pictures, notes, and raw OOXML text that may include grouped
shape or SmartArt text.
"""
from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def ensure_pptx() -> None:
    try:
        import pptx  # noqa: F401
    except ImportError:
        sys.stderr.write("python-pptx is required. Install it in the active environment.\n")
        raise SystemExit(2)


def text_frame_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text:
            text = paragraph.text.strip()
        if text:
            parts.append(text)
    return "\n".join(parts).strip()


def shape_kind(shape) -> str:
    try:
        return str(shape.shape_type).split(".")[-1]
    except Exception:
        return "unknown"


def shape_bbox(shape) -> dict[str, int | None]:
    out: dict[str, int | None] = {}
    for attr in ("left", "top", "width", "height"):
        try:
            out[attr] = int(getattr(shape, attr))
        except Exception:
            out[attr] = None
    return out


def extract_tables(shape) -> list[list[list[str]]]:
    if not getattr(shape, "has_table", False):
        return []
    table = shape.table
    rows: list[list[str]] = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return [rows]


def extract_chart(shape) -> dict | None:
    if not getattr(shape, "has_chart", False):
        return None
    chart = shape.chart
    out: dict = {
        "title": "",
        "categories": [],
        "series": [],
        "axes": [],
        "legend": None,
    }
    try:
        if chart.has_title:
            out["title"] = chart.chart_title.text_frame.text.strip()
    except Exception:
        pass
    try:
        out["legend"] = bool(chart.has_legend)
    except Exception:
        pass
    try:
        plot = chart.plots[0]
        out["categories"] = [str(c) for c in plot.categories]
    except Exception:
        pass
    try:
        for series in chart.series:
            values = []
            try:
                values = [float(v) if isinstance(v, (int, float)) else v for v in series.values]
            except Exception:
                pass
            out["series"].append({"name": str(series.name), "values": values})
    except Exception:
        pass
    try:
        for axis_name in ("category_axis", "value_axis"):
            axis = getattr(chart, axis_name, None)
            if axis is None:
                continue
            title = ""
            try:
                if axis.has_title:
                    title = axis.axis_title.text_frame.text.strip()
            except Exception:
                pass
            out["axes"].append({"axis": axis_name, "title": title})
    except Exception:
        pass
    return out


def extract_shapes(slide) -> list[dict]:
    items: list[dict] = []
    for idx, shape in enumerate(slide.shapes, start=1):
        item: dict = {
            "shape_index": idx,
            "name": getattr(shape, "name", ""),
            "kind": shape_kind(shape),
            "bbox_emu": shape_bbox(shape),
        }
        text = text_frame_text(shape)
        if text:
            item["text"] = text
        tables = extract_tables(shape)
        if tables:
            item["tables"] = tables
        chart = extract_chart(shape)
        if chart:
            item["chart"] = chart
        if "PICTURE" in item["kind"].upper() or "MEDIA" in item["kind"].upper():
            item["visual_note"] = "Picture or media object. Inspect rendered slide and OCR if text is embedded."
        items.append(item)
    return items


def extract_notes(slide) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        return ""


def raw_slide_texts(deck_path: Path) -> dict[int, list[str]]:
    """Collect all DrawingML text nodes from slide XML files."""
    out: dict[int, list[str]] = {}
    with zipfile.ZipFile(deck_path) as zf:
        names = sorted(
            (name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")),
            key=lambda n: int(Path(n).stem.replace("slide", "")),
        )
        for name in names:
            slide_num = int(Path(name).stem.replace("slide", ""))
            root = ET.fromstring(zf.read(name))
            texts = [node.text.strip() for node in root.findall(".//a:t", NS) if node.text and node.text.strip()]
            out[slide_num] = texts
    return out


def extract(deck_path: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(deck_path))
    raw_text = raw_slide_texts(deck_path)
    slides: list[dict] = []
    for idx, slide in enumerate(prs.slides, start=1):
        shapes = extract_shapes(slide)
        title = ""
        for item in shapes:
            if item.get("text"):
                title = str(item["text"]).splitlines()[0]
                break
        shape_texts = []
        for item in shapes:
            if "text" in item:
                shape_texts.extend(str(item["text"]).splitlines())
        raw_only = [text for text in raw_text.get(idx, []) if text not in shape_texts]
        slides.append(
            {
                "slide": idx,
                "title": title,
                "shapes": shapes,
                "raw_ooxml_text": raw_text.get(idx, []),
                "raw_ooxml_text_not_in_shapes": raw_only,
                "existing_notes": extract_notes(slide),
            }
        )
    return {"deck": str(deck_path), "slide_count": len(slides), "slides": slides}


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description="Extract structured .pptx slide evidence.")
    parser.add_argument("deck", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args(argv)

    if not args.deck.exists():
        sys.stderr.write(f"File not found: {args.deck}\n")
        return 1
    ensure_pptx()
    data = extract(args.deck)
    text = json.dumps(data, ensure_ascii=False, indent=2)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(text + "\n", encoding="utf-8")
    else:
        print(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
